"""
Modify the Dissertation Proposal Presentation:
1. Change title to "Dissertation Proposal Defense"
2. Add committee members slide after slide 1
3. Insert chart images into slides 12, 14, 15 (adjusted for new committee slide)
4. Add references slide at the end

Input:  Castro_Puello_Eduardo_Dissertation_Proposal_Presentation_20260508.pptx
Output: Dissertation_Proposal_Defense_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import copy
from lxml import etree

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_PPTX = r'C:\Users\edcastr\OneDrive - Microsoft\Documents\NSU\Summer 2026\Castro_Puello_Eduardo_Dissertation_Proposal_Presentation_20260508.pptx'
OUTPUT_PPTX = os.path.join(SCRIPT_DIR, 'Dissertation_Proposal_Defense_Presentation.pptx')

CHART_DIR = SCRIPT_DIR
CHART_SLIDE12 = os.path.join(CHART_DIR, 'slide12_accuracy_ft_vs_oob.png')
CHART_SLIDE14 = os.path.join(CHART_DIR, 'slide14_blackbox_asr_ft_vs_oob.png')
CHART_SLIDE15 = os.path.join(CHART_DIR, 'slide15_whitebox_asr_ft_vs_oob.png')


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a formatted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def main():
    prs = Presentation(INPUT_PPTX)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ══════════════════════════════════════════════════════
    # 1. MODIFY TITLE (Slide 1)
    # ══════════════════════════════════════════════════════
    print("1. Modifying title slide...")
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if 'Dissertation Proposal' in para.text and 'Defense' not in para.text:
                    for run in para.runs:
                        if 'Dissertation Proposal' in run.text:
                            run.text = run.text.replace('Dissertation Proposal',
                                                       'Dissertation Proposal\nDefense')
                    break

    # ══════════════════════════════════════════════════════
    # 2. ADD COMMITTEE SLIDE (after slide 1)
    # ══════════════════════════════════════════════════════
    print("2. Adding committee slide...")
    # Find blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            blank_layout = layout
            break
    if not blank_layout:
        blank_layout = prs.slide_layouts[-1]

    committee_slide = prs.slides.add_slide(blank_layout)

    # Title
    add_textbox(committee_slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                'Dissertation Committee', font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    # Committee members
    members = [
        ("Committee Chair", "Dr. Wei Li"),
        ("Committee Member", "Dr. Lingwei Wang"),
        ("Committee Member", "Dr. Abhinav Kumar"),
        ("Committee Member", "Dr. Yohannes G. Mulissa"),
    ]

    y_start = Inches(1.5)
    for i, (role, name) in enumerate(members):
        y = y_start + Inches(i * 1.2)
        add_textbox(committee_slide, Inches(2), y, Inches(6), Inches(0.35),
                    role, font_size=13, bold=False,
                    color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)
        add_textbox(committee_slide, Inches(2), y + Inches(0.32), Inches(6), Inches(0.5),
                    name, font_size=22, bold=True,
                    color=RGBColor(0x2C, 0x3E, 0x50), alignment=PP_ALIGN.CENTER)

    # Affiliation
    add_textbox(committee_slide, Inches(1), Inches(6.3), Inches(8), Inches(0.4),
                'College of Computing, AI and Cybersecurity\nNova Southeastern University',
                font_size=11, bold=False,
                color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

    # Move committee slide to position 2 (index 1)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_slide_elem = slides_list[-1]
    xml_slides.remove(new_slide_elem)
    xml_slides.insert(1, new_slide_elem)

    # ══════════════════════════════════════════════════════
    # 3. INSERT CHARTS INTO EXISTING SLIDES
    #    After committee insert: original slide N is now N+1
    #    Slide 12 -> index 12 (0-based: 12)
    #    Slide 14 -> index 14 (0-based: 14)
    #    Slide 15 -> index 15 (0-based: 15)
    # ══════════════════════════════════════════════════════
    print("3. Adding charts to slides...")

    # Slide 13 (was slide 12): Task accuracy chart
    slide12 = prs.slides[12]  # 0-indexed, position 13
    slide12.shapes.add_picture(CHART_SLIDE12,
                               Inches(0.2), Inches(3.8),
                               Inches(9.6), Inches(3.5))
    print(f"   Added accuracy chart to slide 13")

    # Slide 15 (was slide 14): Black-box ASR chart
    slide14 = prs.slides[14]  # 0-indexed, position 15
    slide14.shapes.add_picture(CHART_SLIDE14,
                               Inches(0.1), Inches(3.6),
                               Inches(9.8), Inches(3.7))
    print(f"   Added black-box ASR chart to slide 15")

    # Slide 16 (was slide 15): White-box ASR chart
    slide15 = prs.slides[15]  # 0-indexed, position 16
    slide15.shapes.add_picture(CHART_SLIDE15,
                               Inches(0.1), Inches(3.5),
                               Inches(9.8), Inches(3.8))
    print(f"   Added white-box ASR chart to slide 16")

    # ══════════════════════════════════════════════════════
    # 4. ADD REFERENCES SLIDE (at the end)
    # ══════════════════════════════════════════════════════
    print("4. Adding references slide...")
    ref_slide = prs.slides.add_slide(blank_layout)

    add_textbox(ref_slide, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6),
                'Key References', font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    references = [
        "[1] Castro-Puello, E., Li, W., Wang, L., Kumar, A., & Mulissa, Y. G. (2026). \"Small Language Model Jailbreak Defender Framework.\" SoutheastCon 2026, IEEE.",
        "[2] Yi, S., et al. (2024). \"Jailbreak attacks and defenses against large language models: A survey.\" arXiv:2407.04295.",
        "[3] Hu, E. J., et al. (2022). \"LoRA: Low-rank adaptation of large language models.\" ICLR.",
        "[4] Mazeika, M., et al. (2024). \"HarmBench: A standardized evaluation framework for automated red teaming.\" arXiv:2402.04249.",
        "[5] Zou, A., et al. (2023). \"Universal and transferable adversarial attacks on aligned language models.\" arXiv:2307.15043.",
        "[6] Zhu, S., et al. (2023). \"AutoDAN: Interpretable gradient-based adversarial attacks on LLMs.\" arXiv:2310.15140.",
        "[7] Qi, X., et al. (2023). \"Fine-tuning aligned language models compromises safety.\" arXiv:2310.03693.",
        "[8] Lermen, S., et al. (2023). \"LoRA fine-tuning efficiently undoes safety training in Llama 2-Chat 70B.\" arXiv.",
        "[9] Dettmers, T., et al. (2023). \"QLoRA: Efficient finetuning of quantized LLMs.\" NeurIPS.",
        "[10] Zhang, W., et al. (2025). \"Can Small Language Models Reliably Resist Jailbreak Attacks?\" arXiv:2503.06519.",
        "[11] Zeng, Y., et al. (2024). \"How Johnny can persuade LLMs to jailbreak them.\" arXiv:2401.06373.",
        "[12] Geisler, S., et al. (2024). \"Attacking LLMs with projected gradient descent.\" arXiv:2402.09154.",
    ]

    txBox = ref_slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(6.4))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(references):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(4)

    # ══════════════════════════════════════════════════════
    # SAVE
    # ══════════════════════════════════════════════════════
    prs.save(OUTPUT_PPTX)
    print(f"\nPresentation saved: {OUTPUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
